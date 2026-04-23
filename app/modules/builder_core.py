"""
builder_core.py — PPTX Builder Core Functions
============================================
Layout: LEFT (chart + table, 60%) | RIGHT (AI insight, 40%)
        If no table → LEFT (3 image placeholders) | RIGHT (AI insight)
All text output in English.

All visual settings are driven by design_template.json — no need to edit this file
for colors, sizes, positions, or enable/disable toggles.
"""

import json
import os
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

# ── Design config (loaded once, overridden by design_template.json) ───────────

_D: dict = {}   # full design dict, set by _apply_design()

def _hex(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _color(key: str) -> RGBColor:
    """Resolve a color key (e.g. 'bg_dark') or a hex string (e.g. '#1E2761')."""
    val = _D.get("colors", {}).get(key, key)
    return _hex(val)

def _cfg(*path, default=None):
    """Safe nested get: _cfg('title_slide', 'accent_bar', 'enabled')"""
    node = _D
    for k in path:
        if not isinstance(node, dict):
            return default
        node = node.get(k)
        if node is None:
            return default
    return node

def _inches(*path, default=0.0):
    v = _cfg(*path, default=default)
    return Inches(float(v))

def _pt(*path, default=12):
    v = _cfg(*path, default=default)
    return Pt(float(v))

def _align_map(s: str) -> PP_ALIGN:
    return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
        str(s).lower(), PP_ALIGN.LEFT
    )

# ── Derived layout globals (recalculated after design load) ───────────────────

SLIDE_W: Emu = Inches(13.33)
SLIDE_H: Emu = Inches(7.5)
HEADER_H: Emu = Inches(1.15)
CONTENT_TOP: Emu = Inches(1.25)
CONTENT_BOTTOM: Emu = Inches(0.35)
MARGIN_L: Emu = Inches(0.45)
MARGIN_R: Emu = Inches(0.45)
SPLIT: float = 0.60
GAP: Emu = Inches(0.20)

LEFT_W: Emu = Inches(1)
RIGHT_W: Emu = Inches(1)
LEFT_X: Emu = Inches(0.45)
RIGHT_X: Emu = Inches(1)
CONTENT_H: Emu = Inches(1)

CHART_PALETTE: list = []
FONT_MAIN: str = "Calibri"
FONT_HEADING: str = "Calibri"

CHART_TYPE_MAP = {
    "bar":         XL_CHART_TYPE.BAR_CLUSTERED,
    "grouped_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line":        XL_CHART_TYPE.LINE,
    "pie":         XL_CHART_TYPE.PIE,
}

_using_template: bool = False

def _recalc_layout():
    global SLIDE_W, SLIDE_H, HEADER_H, CONTENT_TOP, CONTENT_BOTTOM
    global MARGIN_L, MARGIN_R, SPLIT, LEFT_W, RIGHT_W, LEFT_X, RIGHT_X, CONTENT_H

    SLIDE_W  = _inches("slide", "width_inches",  default=13.33)
    SLIDE_H  = _inches("slide", "height_inches", default=7.5)
    MARGIN_L = _inches("slide", "margin_left_inches",  default=0.45)
    MARGIN_R = _inches("slide", "margin_right_inches", default=0.45)

    HEADER_H      = _inches("header_band", "height_inches",             default=1.15)
    CONTENT_TOP   = _inches("content_slide", "content_top_inches",      default=1.25)
    CONTENT_BOTTOM= _inches("content_slide", "content_bottom_margin_inches", default=0.35)

    SPLIT  = float(_cfg("split_layout", "split_ratio", default=0.60))
    GAP_   = Inches(0.20)

    LEFT_X  = MARGIN_L
    LEFT_W  = Emu(int(SLIDE_W * SPLIT - int(MARGIN_L) - int(GAP_) // 2))
    RIGHT_W = Emu(int(SLIDE_W * (1 - SPLIT) - int(MARGIN_R) - int(GAP_) // 2))
    RIGHT_X = Emu(int(SLIDE_W * SPLIT + int(GAP_) // 2))
    CONTENT_H = Emu(int(SLIDE_H) - int(CONTENT_TOP) - int(CONTENT_BOTTOM))

def _apply_design(design: dict):
    global _D, FONT_MAIN, FONT_HEADING, CHART_PALETTE
    _D = design

    FONT_MAIN    = _cfg("fonts", "main",    default="Calibri")
    FONT_HEADING = _cfg("fonts", "heading", default=FONT_MAIN)

    palette_raw = _cfg("colors", "chart_palette", default=[])
    CHART_PALETTE[:] = [_hex(h) for h in palette_raw] if palette_raw else [
        RGBColor(0x1E, 0x27, 0x61), RGBColor(0x02, 0x80, 0x90),
        RGBColor(0xCA, 0xDC, 0xFC), RGBColor(0xF9, 0xA8, 0x25),
        RGBColor(0xE5, 0x39, 0x35),
    ]

    _recalc_layout()

# ── Utility ───────────────────────────────────────────────────────────────────

def _fmt(value) -> str:
    if isinstance(value, bool):
        return str(value)
    try:
        v = float(value)
        if abs(v) >= 1_000_000_000: return f"{v/1_000_000_000:.1f}B"
        if abs(v) >= 1_000_000:     return f"{v/1_000_000:.1f}M"
        if 0 < abs(v) < 1:          return f"{v:.1%}"
        if abs(v) >= 1_000:         return f"{v:,.0f}"
        return f"{v:g}"
    except (TypeError, ValueError):
        return str(value)

def _unmask_fmt(text: str, mask_map: dict) -> str:
    result = str(text)
    for token, original in mask_map.items():
        if token in result:
            result = result.replace(token, _fmt(original))
    return result

def _add_text(slide, text: str, left, top, width, height,
              bold=False, size_pt=14, color: RGBColor = None,
              align=PP_ALIGN.LEFT, word_wrap=True, italic=False, font=None):
    if color is None:
        color = _color("body")
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.name = font or FONT_MAIN
    return tb

def _to_float(val) -> float:
    try:
        return float(val)
    except (TypeError, ValueError):
        return 0.0

# ── Remove all existing slides from template ──────────────────────────────────

from .builder_utils import _remove_all_slides, _get_layout, _solid_fill, _solid_fill_alpha
from .builder_components import (
    build_header_band, _add_divider,
    build_title_slide,
    build_right_insight_panel,
    build_left_chart, build_left_table,
    build_left_chart_and_table,
    build_left_image, build_single_image_placeholder
)

# ── Bullet slide ──────────────────────────────────────────────────────────────

def build_bullet_slide(prs, raw_slide, slide_plan, mask_map):
    from pptx.oxml.ns import qn
    from lxml import etree

    color_scheme = slide_plan.get("color_scheme", "corporate_blue")
    slide = prs.slides.add_slide(_get_layout(prs, 6))

    if not _using_template and _cfg("content_slide", "bg", "enabled", default=True):
        bg_c = _color(_cfg("content_slide", "bg", "color", default="bg_light"))
        bg   = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        _solid_fill(bg, bg_c)
        bg.line.fill.background()

    title = _unmask_fmt(str(raw_slide.get("title", "")), mask_map)
    build_header_band(slide, title, color_scheme)

    bullets = raw_slide.get("content", [])
    if not isinstance(bullets, list):
        bullets = [str(bullets)]

    dot_char  = _cfg("bullet_slide", "bullet_dot_char",  default="●")
    dot_color = _color(_cfg("bullet_slide", "bullet_dot_color", default="bullet_dot"))
    body_color= _color(_cfg("bullet_slide", "body_color", default="body"))
    bul_fs    = _cfg("bullet_slide", "bullet_size", default=15)
    spacing   = str(int(_cfg("bullet_slide", "line_spacing_pt", default=700)))

    tb = slide.shapes.add_textbox(LEFT_X, CONTENT_TOP, SLIDE_W - MARGIN_L - MARGIN_R, CONTENT_H)
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        bullet_text = _unmask_fmt(str(bullet), mask_map)
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            pPr = para._p.get_or_add_pPr()
            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
            spcPts.set("val", spacing)

        dot = para.add_run()
        dot.text = f"{dot_char} "
        dot.font.size      = Pt(bul_fs)
        dot.font.color.rgb = dot_color
        dot.font.name      = FONT_MAIN

        txt = para.add_run()
        txt.text           = bullet_text
        txt.font.size      = Pt(bul_fs)
        txt.font.color.rgb = body_color
        txt.font.name      = FONT_MAIN

    # Footer insight
    insight = slide_plan.get("insight_text", "")
    if insight and _cfg("bullet_slide", "footer_insight", "enabled", default=True):
        try:
            from ..sanitizer import unmask
            insight = unmask(insight, mask_map)
        except Exception:
            pass
        fi_fs = _cfg("bullet_slide", "footer_insight", "font_size", default=11)
        fi_c  = _color(_cfg("bullet_slide", "footer_insight", "color", default="#555555"))
        _add_text(slide, insight,
                  LEFT_X, SLIDE_H - bc._inches("bullet_slide", "footer_insight", "bottom_offset", default=0.55),
                  SLIDE_W - MARGIN_L - MARGIN_R, bc._inches("bullet_slide", "footer_insight", "height_inches", default=0.45),
                  size_pt=fi_fs, color=fi_c, italic=True)

# ── Split slide ───────────────────────────────────────────────────────────────

def build_split_slide(prs, raw_slide, slide_plan, mask_map, block_rects=None):
    # Dùng local variables thay vì mutate globals — tránh layout leak sang slide kế tiếp
    local_left_x    = LEFT_X
    local_left_w    = LEFT_W
    local_right_x   = RIGHT_X
    local_right_w   = RIGHT_W
    local_cont_top  = CONTENT_TOP
    local_cont_h    = CONTENT_H

    if block_rects and len(block_rects) >= 2:
        r0, r1 = block_rects[0], block_rects[1]
        local_left_x   = Inches(r0.left)
        local_left_w   = Inches(r0.width)
        local_right_x  = Inches(r1.left)
        local_right_w  = Inches(r1.width)
        local_cont_top = Inches(r0.top)
        local_cont_h   = Inches(r0.height)
    elif block_rects and len(block_rects) == 1:
        r0 = block_rects[0]
        local_left_x   = Inches(r0.left)
        local_left_w   = Inches(r0.width)
        local_cont_top = Inches(r0.top)
        local_cont_h   = Inches(r0.height)

    color_scheme   = slide_plan.get("color_scheme", "corporate_blue")
    chart_type_key = slide_plan.get("chart_type", "none")
    insight_text   = slide_plan.get("insight_text", "")

    slide = prs.slides.add_slide(_get_layout(prs, 6))

    if not _using_template and _cfg("content_slide", "bg", "enabled", default=True):
        bg_c = _color(_cfg("content_slide", "bg", "color", default="bg_light"))
        bg   = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        _solid_fill(bg, bg_c)
        bg.line.fill.background()

    title = _unmask_fmt(str(raw_slide.get("title", "")), mask_map)
    build_header_band(slide, title, color_scheme)

    # Tạm thời patch globals để builder_components dùng đúng local layout
    _saved = (LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CONTENT_TOP, CONTENT_H)
    globals().update({
        'LEFT_X': local_left_x, 'LEFT_W': local_left_w,
        'RIGHT_X': local_right_x, 'RIGHT_W': local_right_w,
        'CONTENT_TOP': local_cont_top, 'CONTENT_H': local_cont_h,
    })
    try:
        if slide_plan.get("show_insight", True):
            _add_divider(slide)

        content  = raw_slide.get("content", {})
        has_table = isinstance(content, dict) and bool(content.get("rows"))

        if has_table:
            if chart_type_key and chart_type_key != "none":
                build_left_chart_and_table(slide, raw_slide, mask_map, chart_type_key, color_scheme)
            else:
                build_left_table(slide, raw_slide, mask_map, color_scheme)
        else:
            images = raw_slide.get('images', [])
            if images and images[0]:
                build_left_image(slide, images[0], title)
            else:
                build_single_image_placeholder(slide, title)

        if slide_plan.get("show_insight", True):
            build_right_insight_panel(slide, insight_text, mask_map)
    finally:
        # Restore globals về giá trị gốc sau mỗi slide
        globals().update({
            'LEFT_X': _saved[0], 'LEFT_W': _saved[1],
            'RIGHT_X': _saved[2], 'RIGHT_W': _saved[3],
            'CONTENT_TOP': _saved[4], 'CONTENT_H': _saved[5],
        })

# ── Grid slide (table2 / table3) ─────────────────────────────────────────────

def build_grid_slide(prs, raw_slide, slide_plan, mask_map, block_rects):
    """
    Render N-column layout dùng block_rects từ layout_engine.
    Cột 0: content chính (table + chart nếu có)
    Cột 1+: extra_tables[i-1] nếu tồn tại, ngược lại dùng insight panel
    """
    color_scheme   = slide_plan.get("color_scheme", "corporate_blue")
    chart_type_key = slide_plan.get("chart_type", "none")
    insight_text   = slide_plan.get("insight_text", "")
    extra_tables   = raw_slide.get("extra_tables", [])

    slide = prs.slides.add_slide(_get_layout(prs, 6))

    if not _using_template and _cfg("content_slide", "bg", "enabled", default=True):
        bg_c = _color(_cfg("content_slide", "bg", "color", default="bg_light"))
        bg   = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        _solid_fill(bg, bg_c)
        bg.line.fill.background()

    title = _unmask_fmt(str(raw_slide.get("title", "")), mask_map)
    build_header_band(slide, title, color_scheme)
    if slide_plan.get("show_insight", True) or len(extra_tables) > 0:
        _add_divider(slide)

    _saved = (LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CONTENT_TOP, CONTENT_H)

    for col_i, rect in enumerate(block_rects):
        col_left = Inches(rect.left)
        col_w    = Inches(rect.width)
        col_top  = Inches(rect.top)
        col_h    = Inches(rect.height)

        # Patch globals cho builder_components
        globals().update({
            'LEFT_X': col_left, 'LEFT_W': col_w,
            'RIGHT_X': col_left, 'RIGHT_W': col_w,
            'CONTENT_TOP': col_top, 'CONTENT_H': col_h,
        })

        try:
            if col_i == 0:
                # Cột đầu: content chính
                content = raw_slide.get("content", {})
                has_table = isinstance(content, dict) and bool(content.get("rows"))
                if has_table:
                    if chart_type_key and chart_type_key != "none":
                        build_left_chart_and_table(slide, raw_slide, mask_map, chart_type_key, color_scheme)
                    else:
                        build_left_table(slide, raw_slide, mask_map, color_scheme)
                else:
                    images = raw_slide.get('images', [])
                    if images and images[0]:
                        build_left_image(slide, images[0], title)
                    else:
                        build_single_image_placeholder(slide, title)
            elif col_i - 1 < len(extra_tables):
                # Cột tiếp theo: extra_table tương ứng
                fake_slide = dict(raw_slide)
                fake_slide["content"] = extra_tables[col_i - 1]
                if chart_type_key and chart_type_key != "none":
                    build_left_chart_and_table(slide, fake_slide, mask_map, chart_type_key, color_scheme)
                else:
                    build_left_table(slide, fake_slide, mask_map, color_scheme)
            else:
                # Cột cuối không có data: insight panel
                if slide_plan.get("show_insight", True):
                    build_right_insight_panel(slide, insight_text, mask_map)
        finally:
            pass  # Restore sau vòng lặp

    # Restore globals
    globals().update({
        'LEFT_X': _saved[0], 'LEFT_W': _saved[1],
        'RIGHT_X': _saved[2], 'RIGHT_W': _saved[3],
        'CONTENT_TOP': _saved[4], 'CONTENT_H': _saved[5],
    })


# ── Entry point ───────────────────────────────────────────────────────────────

def build_dynamic_slide(prs, raw_slide, slide_plan, mask_map, block_rects, layout_cfg):
    color_scheme = slide_plan.get("color_scheme", "corporate_blue")
    slide = prs.slides.add_slide(_get_layout(prs, 6))

    if not _using_template and _cfg("content_slide", "bg", "enabled", default=True):
        bg_c = _color(_cfg("content_slide", "bg", "color", default="bg_light"))
        bg   = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
        _solid_fill(bg, bg_c)
        bg.line.fill.background()

    title = _unmask_fmt(str(raw_slide.get("title", "")), mask_map)
    build_header_band(slide, title, color_scheme)

    components = layout_cfg.get('components', [])
    active_comps = [c for c in components if c.get('active', True)]
    if any(c.get('type') == 'insight' for c in active_comps):
        _add_divider(slide)

    chart_type_key = slide_plan.get("chart_type", "bar")
    if chart_type_key == "none":
        chart_type_key = "bar" # Fallback if user wants chart but AI said none

    insight_text = slide_plan.get("insight_text", "")
    _saved = (LEFT_X, LEFT_W, RIGHT_X, RIGHT_W, CONTENT_TOP, CONTENT_H)

    for i, comp in enumerate(active_comps):
        if not block_rects or i >= len(block_rects):
            break
            
        rect = block_rects[i]
        col_left = Inches(rect.left)
        col_w    = Inches(rect.width)
        col_top  = Inches(rect.top)
        col_h    = Inches(rect.height)

        globals().update({
            'LEFT_X': col_left, 'LEFT_W': col_w,
            'RIGHT_X': col_left, 'RIGHT_W': col_w,
            'CONTENT_TOP': col_top, 'CONTENT_H': col_h,
        })

        try:
            c_type = comp.get('type')
            c_id = str(comp.get('id', ''))
            
            # Determine which table data to use based on id like 'tbl0', 'tbl1', 'cht0', 'cht1'
            table_idx = 0
            if c_id.startswith('tbl'):
                table_idx = int(c_id.replace('tbl', '')) if c_id.replace('tbl', '').isdigit() else 0
            elif c_id.startswith('cht'):
                table_idx = int(c_id.replace('cht', '')) if c_id.replace('cht', '').isdigit() else 0

            data_slide = dict(raw_slide)
            if table_idx > 0:
                extra_tables = raw_slide.get("extra_tables", [])
                if table_idx - 1 < len(extra_tables):
                    data_slide["content"] = extra_tables[table_idx - 1]

            if c_type == 'table':
                build_left_table(slide, data_slide, mask_map, color_scheme)
            elif c_type == 'chart':
                build_left_chart(slide, data_slide, mask_map, chart_type_key)
            elif c_type == 'content':
                content_data = data_slide.get("content", [])
                if isinstance(content_data, list) and content_data and isinstance(content_data[0], str):
                    bullets = content_data
                    dot_char  = _cfg("bullet_slide", "bullet_dot_char",  default="●")
                    dot_color = _color(_cfg("bullet_slide", "bullet_dot_color", default="bullet_dot"))
                    body_color= _color(_cfg("bullet_slide", "body_color", default="body"))
                    bul_fs    = _cfg("bullet_slide", "bullet_size", default=15)
                    spacing   = str(int(_cfg("bullet_slide", "line_spacing_pt", default=700)))

                    tb = slide.shapes.add_textbox(col_left, col_top, col_w, col_h)
                    tf = tb.text_frame
                    tf.word_wrap = True

                    from pptx.oxml.ns import qn
                    from lxml import etree
                    for bi, bullet in enumerate(bullets):
                        bullet_text = _unmask_fmt(str(bullet), mask_map)
                        para = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                        if bi > 0:
                            pPr = para._p.get_or_add_pPr()
                            spcBef = etree.SubElement(pPr, qn("a:spcBef"))
                            spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
                            spcPts.set("val", spacing)

                        dot = para.add_run()
                        dot.text = f"{dot_char} "
                        dot.font.size      = Pt(bul_fs)
                        dot.font.color.rgb = dot_color
                        dot.font.name      = FONT_MAIN

                        txt = para.add_run()
                        txt.text           = bullet_text
                        txt.font.size      = Pt(bul_fs)
                        txt.font.color.rgb = body_color
                        txt.font.name      = FONT_MAIN
                else:
                    images = data_slide.get('images', [])
                    if images and images[0]:
                        build_left_image(slide, images[0], title)
                    else:
                        build_single_image_placeholder(slide, title)
            elif c_type == 'image':
                # determine image index
                img_idx = int(c_id.replace('img', '')) if c_id.replace('img', '').isdigit() else 0
                
                # try to get url from layout_cfg if present (AI generated)
                image_url = None
                if layout_cfg and 'imageUrls' in layout_cfg and layout_cfg['imageUrls']:
                    if img_idx < len(layout_cfg['imageUrls']) and layout_cfg['imageUrls'][img_idx]:
                        image_url = layout_cfg['imageUrls'][img_idx]
                
                # if not AI generated, use image from raw docx extraction
                if not image_url:
                    images = data_slide.get('images', [])
                    if img_idx < len(images) and images[img_idx]:
                        image_url = images[img_idx]
                
                if image_url:
                    build_left_image(slide, image_url, title)
                else:
                    build_single_image_placeholder(slide, title, img_idx + 1)
            elif c_type == 'insight':
                build_right_insight_panel(slide, insight_text, mask_map)
        except Exception as e:
            print(f"Error building component {c_type}: {e}")
        finally:
            pass

    globals().update({
        'LEFT_X': _saved[0], 'LEFT_W': _saved[1],
        'RIGHT_X': _saved[2], 'RIGHT_W': _saved[3],
        'CONTENT_TOP': _saved[4], 'CONTENT_H': _saved[5],
    })

def build_slide(prs, raw_slide, slide_plan, mask_map, scale_map=None, block_rects=None, layout_cfg=None):
    slide_type = str(raw_slide.get("type", "")).lower()
    layout     = str(slide_plan.get("layout", "one_col")).lower()

    if slide_type == "title" or layout == "title_only":
        build_title_slide(prs, raw_slide, mask_map)
        return

    if layout_cfg and layout_cfg.get('components'):
        build_dynamic_slide(prs, raw_slide, slide_plan, mask_map, block_rects, layout_cfg)
        return

    if layout in ("table2", "table3") and block_rects and len(block_rects) >= 2:
        build_grid_slide(prs, raw_slide, slide_plan, mask_map, block_rects)
        return

    build_split_slide(prs, raw_slide, slide_plan, mask_map, block_rects=block_rects)

def build(layout_plan: dict, mask_map: dict, raw_data: dict,
          scale_map: dict = None,
          template_path: str = None,
          output_path: str = "output.pptx",
          design: dict = None,
          block_rects: list = None,
          per_slide_rects: list = None,
          layouts_cfg: list = None) -> None:

    if design is not None:
        _apply_design(design)
    else:
        # Load defaults so _D is never empty
        _apply_design({})

    global _using_template
    if template_path and os.path.exists(template_path):
        _using_template = True
        prs = Presentation(template_path)
        _remove_all_slides(prs)
    else:
        _using_template = False
        prs = Presentation()
        prs.slide_width  = SLIDE_W
        prs.slide_height = SLIDE_H

    # Cover slide
    pres_title = raw_data.get("title") or raw_data.get("presentation_title", "")
    if pres_title:
        cover = {"type": "title", "title": pres_title}
        build_title_slide(prs, cover, mask_map)

    plan_slides = layout_plan.get("slides", [])
    raw_slides  = raw_data.get("slides", [])

    for i, slide_plan in enumerate(plan_slides):
        idx = slide_plan.get("index", 0)
        if idx < len(raw_slides):
            raw_slide = raw_slides[idx]
        elif raw_slides:
            raw_slide = raw_slides[-1]
        else:
            continue
        slide_rects = (per_slide_rects[i]
                       if per_slide_rects and i < len(per_slide_rects)
                       else block_rects)
        slide_layout_cfg = layouts_cfg[i] if layouts_cfg and i < len(layouts_cfg) else None
        build_slide(prs, raw_slide, slide_plan, mask_map, scale_map, block_rects=slide_rects, layout_cfg=slide_layout_cfg)

    prs.save(output_path)
    print(f"\n✅ Saved: {output_path}")
