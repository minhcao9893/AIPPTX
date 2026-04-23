"""
builder_components.py — Reusable slide-building blocks.

Imports from builder_core are done lazily (inside each function) to avoid
circular-import issues: builder_core imports this module at the bottom, so a
top-level 'from .builder_core import *' would run before builder_core has
finished initializing its globals (_D, _cfg, FONT_HEADING, …).
"""

from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData


def _bc():
    """Lazy accessor for builder_core namespace (avoids circular import)."""
    from . import builder_core as _m
    return _m


# ── Header band ───────────────────────────────────────────────────────────────

def build_header_band(slide, title_text: str, color_scheme: str = "corporate_blue"):
    bc = _bc()
    if not bc._cfg("header_band", "enabled", default=True):
        return

    hdr_color = bc._color(bc._cfg("header_band", "color", default="bg_dark"))
    h         = bc.HEADER_H
    header    = slide.shapes.add_shape(1, 0, 0, bc.SLIDE_W, h)

    # Alpha transparency: 0 = opaque, 100 = fully transparent
    alpha_pct = bc._cfg('transparency', 'header_band_alpha', default=0)
    try:
        alpha_pct = int(alpha_pct)
    except (TypeError, ValueError):
        alpha_pct = 0
    
    # Always use _solid_fill_alpha to handle transparency correctly
    bc._solid_fill_alpha(header, hdr_color, alpha_pct)
    header.line.fill.background()

    tx = bc._inches("header_band", "title_x_inches",      default=0.45)
    ty = bc._inches("header_band", "title_y_inches",      default=0.16)
    tw = bc._inches("header_band", "title_width_inches",  default=12.4)
    th = bc._inches("header_band", "title_height_inches", default=0.85)
    fs = bc._cfg("fonts", "header_band_size", default=28)

    # Determine title color: use white if background is dark/solid, 
    # otherwise use the theme's heading color (dark).
    auto_text = bc._cfg("header_band", "auto_text_color", default=True)
    if auto_text:
        # If band is very transparent (> 50%), use heading color (dark)
        if alpha_pct > 50:
            title_color = bc._color("heading")
        else:
            title_color = bc._color("#FFFFFF")
    else:
        title_color = bc._color(bc._cfg("header_band", "text_color", default="#FFFFFF"))

    bc._add_text(slide, title_text, tx, ty, tw, th,
                 bold=True, size_pt=fs,
                 color=title_color,
                 font=bc.FONT_HEADING)


# ── Divider line ──────────────────────────────────────────────────────────────

def _add_divider(slide):
    bc = _bc()
    if not bc._cfg("divider", "enabled", default=True):
        return
    div_color = bc._color(bc._cfg("divider", "color", default="divider_line"))
    div_w_raw = bc._cfg("divider", "width_inches", default=0.04)
    div_w     = Inches(float(div_w_raw))

    div_x = int(bc.SLIDE_W * bc.SPLIT)
    div_y = int(bc.CONTENT_TOP)
    div_h = int(bc.CONTENT_H)

    line = slide.shapes.add_shape(
        1,
        div_x - int(div_w) // 2, div_y,
        int(div_w), div_h
    )
    bc._solid_fill(line, div_color)
    line.line.fill.background()


# ── Title slide ───────────────────────────────────────────────────────────────

def build_title_slide(prs, raw_slide, mask_map):
    bc = _bc()
    if not bc._cfg("title_slide", "enabled", default=True):
        return

    from .builder_utils import _get_layout
    slide = prs.slides.add_slide(_get_layout(prs, 0))

    # Background
    if not bc._using_template and bc._cfg("title_slide", "bg", "enabled", default=True):
        bg_color = bc._color(bc._cfg("title_slide", "bg", "color", default="bg_dark"))
        bg = slide.shapes.add_shape(1, 0, 0, bc.SLIDE_W, bc.SLIDE_H)
        bc._solid_fill(bg, bg_color)
        bg.line.fill.background()

    # Accent bar
    if bc._cfg("title_slide", "accent_bar", "enabled", default=True):
        bar_w = bc._inches("title_slide", "accent_bar", "width_inches", default=0.18)
        bar_c = bc._color(bc._cfg("title_slide", "accent_bar", "color", default="bullet_dot"))
        bar   = slide.shapes.add_shape(1, 0, 0, bar_w, bc.SLIDE_H)
        bc._solid_fill(bar, bar_c)
        bar.line.fill.background()

    # Title
    title_text = bc._unmask_fmt(str(raw_slide.get("title", "")), mask_map)
    ty = bc._inches("title_slide", "title", "y_inches",      default=2.2)
    th = bc._inches("title_slide", "title", "height_inches", default=2.0)
    fs = bc._cfg("title_slide", "title", "font_size", default=40)
    tc = bc._color(bc._cfg("title_slide", "title", "color", default="#FFFFFF"))
    al = bc._align_map(bc._cfg("title_slide", "title", "align", default="center"))

    bc._add_text(slide, title_text,
                 bc._inches("title_slide", "title", "x_inches", default=0.55), ty, 
                 bc._inches("title_slide", "title", "width_inches", default=12.33), th,
                 bold=True, size_pt=fs, color=tc, align=al, font=bc.FONT_HEADING)

    # Subtitle
    subtitle = raw_slide.get("subtitle", "")
    if subtitle and bc._cfg("title_slide", "subtitle", "enabled", default=True):
        sub_text = bc._unmask_fmt(str(subtitle), mask_map)
        sy = bc._inches("title_slide", "subtitle", "y_inches",      default=4.5)
        sh = bc._inches("title_slide", "subtitle", "height_inches", default=1.0)
        sf = bc._cfg("title_slide", "subtitle", "font_size", default=18)
        sc = bc._color(bc._cfg("title_slide", "subtitle", "color", default="accent"))
        sa = bc._align_map(bc._cfg("title_slide", "subtitle", "align", default="center"))

        bc._add_text(slide, sub_text,
                     bc._inches("title_slide", "subtitle", "x_inches", default=0.55), sy, 
                     bc._inches("title_slide", "subtitle", "width_inches", default=12.33), sh,
                     size_pt=sf, color=sc, align=sa)


# ── RIGHT PANEL: AI Insight ───────────────────────────────────────────────────

def build_right_insight_panel(slide, insight_text: str, mask_map: dict):
    bc = _bc()
    if not bc._cfg("right_panel", "enabled", default=True):
        return

    if not insight_text:
        insight_text = "No insight available."
    try:
        from ..sanitizer import unmask
        text = unmask(insight_text, mask_map)
    except Exception:
        text = insight_text

    pad   = bc._inches("right_panel", "padding_inches", default=0.15)
    bg_c  = bc._color(bc._cfg("right_panel", "bg_color",     default="insight_bg"))
    brd_c = bc._color(bc._cfg("right_panel", "border_color", default="insight_border"))

    panel_bg = slide.shapes.add_shape(
        1, Emu(int(bc.RIGHT_X)), Emu(int(bc.CONTENT_TOP)),
        Emu(int(bc.RIGHT_W)), Emu(int(bc.CONTENT_H))
    )
    panel_bg.fill.solid()
    panel_bg.fill.fore_color.rgb = bg_c
    panel_bg.line.color.rgb = brd_c

    label_y_offset = Inches(0.12)

    # Label
    if bc._cfg("right_panel", "label", "enabled", default=True):
        label_text = bc._cfg("right_panel", "label", "text",      default="AI Insight")
        label_fs   = bc._cfg("right_panel", "label", "font_size", default=13)
        label_c    = bc._color(bc._cfg("right_panel", "label", "color", default="insight_label"))

        bc._add_text(slide, label_text,
                     bc.RIGHT_X + pad,
                     bc.CONTENT_TOP + label_y_offset,
                     bc.RIGHT_W - pad * 2, bc._inches("right_panel", "label", "height_inches", default=0.35),
                     bold=True, size_pt=label_fs, color=label_c, font=bc.FONT_HEADING)

        # Divider under label
        if bc._cfg("right_panel", "label_divider", "enabled", default=True):
            div_c = bc._color(bc._cfg("right_panel", "label_divider", "color",
                                      default="divider_line"))
            div   = slide.shapes.add_shape(
                1,
                bc.RIGHT_X + pad,
                bc.CONTENT_TOP + bc._inches("right_panel", "label_divider", "top_offset", default=0.50),
                bc.RIGHT_W - pad * 2,
                bc._inches("right_panel", "label_divider", "height_inches", default=0.03)
            )
            bc._solid_fill(div, div_c)
            div.line.fill.background()

    # Body text
    body_fs = bc._cfg("right_panel", "body", "font_size", default=13)
    body_c  = bc._color(bc._cfg("right_panel", "body", "color", default="insight_body"))

    bc._add_text(slide, text,
                 bc.RIGHT_X + pad,
                 bc.CONTENT_TOP + Inches(0.62),
                 bc.RIGHT_W - pad * 2,
                 bc.CONTENT_H - Inches(0.80),
                 size_pt=body_fs, color=body_c, word_wrap=True)


# ── LEFT PANEL: Chart ─────────────────────────────────────────────────────────

def build_left_chart(slide, raw_slide, mask_map, chart_type_key: str):
    bc = _bc()
    content = raw_slide.get("content", {})
    if not isinstance(content, dict) or not content.get("rows"):
        return False

    xl_type     = bc.CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data  = _build_chart_data(raw_slide, mask_map, chart_type_key)
    chart_shape = slide.shapes.add_chart(
        xl_type, bc.LEFT_X, bc.CONTENT_TOP, bc.LEFT_W, bc.CONTENT_H, chart_data
    )
    chart   = chart_shape.chart
    columns = content.get("columns", [])
    chart.has_legend = bool(bc._cfg("chart", "has_legend", default=True)) and len(columns) > 2
    _apply_chart_colors(chart, max(len(columns) - 1, 1))
    return True


def _build_chart_data(raw_slide, mask_map, chart_type_key: str) -> ChartData:
    bc      = _bc()
    content = raw_slide.get("content", {})
    columns = content.get("columns", [])
    rows    = content.get("rows", [])
    chart_data = ChartData()
    if not rows:
        return chart_data

    if chart_type_key == "pie":
        categories = [
            bc._unmask_fmt(str(r[0]), mask_map) if isinstance(r, list) else str(i)
            for i, r in enumerate(rows)
        ]
        chart_data.categories = categories
        values = [
            bc._to_float(r[1]) if isinstance(r, list) and len(r) >= 2 else 0.0
            for r in rows
        ]
        series_name = (bc._unmask_fmt(str(columns[1]), mask_map)
                       if len(columns) > 1 else "Value")
        chart_data.add_series(series_name, values)
    else:
        label_col = [
            bc._unmask_fmt(str(r[0]), mask_map) if isinstance(r, list) and r else f"Item {i}"
            for i, r in enumerate(rows)
        ]
        chart_data.categories = label_col
        n_series = max(len(columns) - 1, 1)
        for si in range(n_series):
            series_name = (
                bc._unmask_fmt(str(columns[si + 1]), mask_map)
                if si + 1 < len(columns) else f"Series {si+1}"
            )
            values = [
                bc._to_float(row[si + 1])
                if isinstance(row, list) and si + 1 < len(row) else 0.0
                for row in rows
            ]
            chart_data.add_series(series_name, values)

    return chart_data


def _apply_chart_colors(chart, n_series: int):
    """Áp màu palette lên từng series/point và bật data labels."""
    from pptx.oxml.ns import qn
    bc = _bc()
    try:
        xl_type = chart.chart_type
        is_pie  = str(xl_type) in ('PIE', 'DOUGHNUT', 'PIE_EXPLODED')
        show_labels = bool(bc._cfg("chart", "show_data_labels", default=True))

        for i, series in enumerate(chart.series):
            if is_pie:
                try:
                    for j, point in enumerate(series.points):
                        color = bc.CHART_PALETTE[j % len(bc.CHART_PALETTE)]
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = color
                except Exception:
                    color = bc.CHART_PALETTE[i % len(bc.CHART_PALETTE)]
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = color
            else:
                color = bc.CHART_PALETTE[i % len(bc.CHART_PALETTE)]
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color

            if show_labels:
                try:
                    dLbls = series._element.find(qn('c:dLbls'))
                    if dLbls is None:
                        from lxml import etree
                        dLbls = etree.SubElement(series._element, qn('c:dLbls'))
                    showVal = dLbls.find(qn('c:showVal'))
                    if showVal is None:
                        from lxml import etree
                        showVal = etree.SubElement(dLbls, qn('c:showVal'))
                    showVal.set('val', '1')
                    for tag in ['c:showLegendKey', 'c:showCatName',
                                'c:showSerName', 'c:showPercent']:
                        el = dLbls.find(qn(tag))
                        if el is None:
                            from lxml import etree
                            el = etree.SubElement(dLbls, qn(tag))
                        el.set('val', '1' if (tag == 'c:showPercent' and is_pie) else '0')
                except Exception:
                    pass
    except Exception:
        pass


# ── LEFT PANEL: Table ─────────────────────────────────────────────────────────

def build_left_table(slide, raw_slide, mask_map, color_scheme: str,
                     top_offset=None, height_override=None):
    bc      = _bc()
    content = raw_slide.get("content", {})
    columns = content.get("columns", [])
    rows    = content.get("rows", [])
    if not columns or not rows:
        return False

    n_rows = len(rows) + 1
    n_cols = len(columns)
    top    = top_offset    if top_offset    is not None else bc.CONTENT_TOP
    height = height_override if height_override is not None else bc.CONTENT_H

    table_shape = slide.shapes.add_table(n_rows, n_cols, bc.LEFT_X, top, bc.LEFT_W, height)
    table = table_shape.table
    col_w = Emu(int(bc.LEFT_W) // n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    hdr_bg_c = bc._color(bc._cfg("table", "header", "bg_color",   default="table_header"))
    hdr_fc   = bc._color(bc._cfg("table", "header", "color",       default="#FFFFFF"))
    hdr_fs   = bc._cfg("table", "header", "font_size",             default=12)
    body_fs  = bc._cfg("table", "body",   "font_size",             default=11)
    alt_c    = bc._color(bc._cfg("table", "body", "alt_row_color", default="table_row_alt"))

    for ci, col_name in enumerate(columns):
        cell = table.cell(0, ci)
        cell.text = bc._unmask_fmt(str(col_name), mask_map)
        para = cell.text_frame.paragraphs[0]
        run  = para.runs[0] if para.runs else para.add_run()
        run.font.bold      = True
        run.font.size      = Pt(hdr_fs)
        run.font.color.rgb = hdr_fc
        run.font.name      = bc.FONT_HEADING
        cell.fill.solid()
        cell.fill.fore_color.rgb = hdr_bg_c

    for ri, row in enumerate(rows):
        row_data = row if isinstance(row, list) else list(row.values())
        for ci, cell_val in enumerate(row_data[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.text = bc._unmask_fmt(str(cell_val), mask_map)
            para = cell.text_frame.paragraphs[0]
            run  = para.runs[0] if para.runs else para.add_run()
            run.font.size  = Pt(body_fs)
            run.font.name  = bc.FONT_MAIN
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_c

    return True


# ── LEFT PANEL: Chart + Table stacked ────────────────────────────────────────

def build_left_chart_and_table(slide, raw_slide, mask_map,
                                chart_type_key: str, color_scheme: str) -> bool:
    bc      = _bc()
    content = raw_slide.get("content", {})
    if not isinstance(content, dict) or not content.get("rows"):
        return False

    chart_h_pct = bc._cfg("stacked_layout", "chart_height_ratio", default=0.58)
    table_h_pct = bc._cfg("stacked_layout", "table_height_ratio", default=0.38)
    gap_h_pct   = bc._cfg("stacked_layout", "gap_ratio",          default=0.04)

    chart_h = Emu(int(int(bc.CONTENT_H) * chart_h_pct))
    table_h = Emu(int(int(bc.CONTENT_H) * table_h_pct))
    gap_h   = Emu(int(int(bc.CONTENT_H) * gap_h_pct))

    xl_type    = bc.CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = _build_chart_data(raw_slide, mask_map, chart_type_key)
    chart_shape = slide.shapes.add_chart(
        xl_type,
        Emu(int(bc.LEFT_X)), Emu(int(bc.CONTENT_TOP)),
        Emu(int(bc.LEFT_W)), chart_h,
        chart_data
    )
    chart   = chart_shape.chart
    columns = content.get("columns", [])
    chart.has_legend = bool(bc._cfg("chart", "has_legend", default=True)) and len(columns) > 2
    _apply_chart_colors(chart, max(len(columns) - 1, 1))

    table_top = Emu(int(bc.CONTENT_TOP) + int(chart_h) + int(gap_h))
    build_left_table(slide, raw_slide, mask_map, color_scheme,
                     top_offset=table_top, height_override=table_h)
    return True


# ── LEFT PANEL: Image Placeholders ───────────────────────────────────────────

def build_left_image_placeholders(slide, slide_title: str):
    bc = _bc()
    n      = int(bc._cfg("image_placeholder", "count",           default=3))
    fill_c = bc._color(bc._cfg("image_placeholder", "fill_color",   default="img_placeholder_fill"))
    brd_c  = bc._color(bc._cfg("image_placeholder", "border_color", default="img_placeholder_border"))
    lbl_fs = bc._cfg("image_placeholder", "label_font_size",         default=11)
    lbl_c  = bc._color(bc._cfg("image_placeholder", "label_color",   default="#5555AA"))

    pad   = Emu(int(Inches(0.12)))
    box_h = Emu(int((int(bc.CONTENT_H) - int(pad) * (n - 1)) / n))

    for i in range(n):
        bx = Emu(int(bc.LEFT_X))
        by = Emu(int(bc.CONTENT_TOP) + i * (int(box_h) + int(pad)))
        bw = Emu(int(bc.LEFT_W))
        bh = box_h

        box = slide.shapes.add_shape(1, bx, by, bw, bh)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_c
        box.line.color.rgb = brd_c
        box.line.width = Pt(1)

        label = (f"[Image {i+1} — related to: {slide_title}]\n"
                 f"Replace or delete this placeholder")
        bc._add_text(slide, label,
                     bx + Emu(int(Inches(0.15))),
                     Emu(int(by) + int(int(bh) * 0.30)),
                     bw - Emu(int(Inches(0.3))),
                     Emu(int(int(bh) * 0.5)),
                     size_pt=lbl_fs, color=lbl_c,
                     align=PP_ALIGN.CENTER, italic=True)

def build_single_image_placeholder(slide, slide_title: str, img_idx: int = 1):
    bc = _bc()
    fill_c = bc._color(bc._cfg("image_placeholder", "fill_color",   default="img_placeholder_fill"))
    brd_c  = bc._color(bc._cfg("image_placeholder", "border_color", default="img_placeholder_border"))
    lbl_fs = bc._cfg("image_placeholder", "label_font_size",         default=11)
    lbl_c  = bc._color(bc._cfg("image_placeholder", "label_color",   default="#5555AA"))

    bx = Emu(int(bc.LEFT_X))
    by = Emu(int(bc.CONTENT_TOP))
    bw = Emu(int(bc.LEFT_W))
    bh = Emu(int(bc.CONTENT_H))

    box = slide.shapes.add_shape(1, bx, by, bw, bh)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_c
    box.line.color.rgb = brd_c
    box.line.width = Pt(1)

    label = f"[ Hình ảnh {img_idx} ]\nChưa có ảnh, vui lòng tạo ảnh bằng AI hoặc chèn ảnh thủ công."
    bc._add_text(slide, label,
                 bx + Emu(int(Inches(0.15))),
                 Emu(int(by) + int(int(bh) * 0.40)),
                 bw - Emu(int(Inches(0.3))),
                 Emu(int(int(bh) * 0.4)),
                 size_pt=lbl_fs, color=lbl_c,
                 align=PP_ALIGN.CENTER, italic=True)

# ── LEFT PANEL: Real Image ───────────────────────────────────────────────────

def build_left_image(slide, image_source: str, slide_title: str):
    bc = _bc()
    import os, tempfile, urllib.request
    from pathlib import Path

    local_path = None
    tmp_file = None

    if image_source.startswith("http://") or image_source.startswith("https://"):
        try:
            req = urllib.request.Request(image_source, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req) as response:
                img_data = response.read()
            tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            tmp_file.write(img_data)
            tmp_file.close()
            local_path = tmp_file.name
        except Exception as e:
            print(f"Error downloading image {image_source}: {e}")
            build_single_image_placeholder(slide, slide_title)
            return
    else:
        # Strip potential API prefix if sent from UI
        clean_source = image_source.replace("/api/media/", "")
        
        # Check if it's a local filename in app/input/media
        # Note: image_source could be just a filename from docx parser
        # __file__ is app/modules/builder_components.py, so parents[1] is app/
        media_dir = Path(__file__).resolve().parents[1] / "input" / "media"
        
        if (media_dir / clean_source).exists():
            local_path = str(media_dir / clean_source)
        elif os.path.exists(image_source):
            local_path = image_source
        elif os.path.exists(clean_source):
            local_path = clean_source
        else:
            print(f"Local image not found: {clean_source} (searched in {media_dir})")
            build_single_image_placeholder(slide, slide_title)
            return

    try:
        slide.shapes.add_picture(local_path, bc.LEFT_X, bc.CONTENT_TOP, bc.LEFT_W, bc.CONTENT_H)
    except Exception as e:
        print(f"Error adding picture {local_path}: {e}")
        build_single_image_placeholder(slide, slide_title)

    if tmp_file:
        try:
            os.unlink(tmp_file.name)
        except:
            pass
